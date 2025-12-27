import pytesseract
import PyPDF2 
from PIL import Image # Add this import: pip install Pillow
import git  # Ensure you have run 'pip install GitPython'
from git.exc import InvalidGitRepositoryError
from pathlib import Path
import hashlib
from firebase_admin import firestore
from collections import OrderedDict
import time
from langchain_text_splitters import RecursiveCharacterTextSplitter
from flask import request, jsonify
import fitz
import io
import re
import cv2
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

TXT_OUTPUT_DIR = Path("converted_txt_projects") # Main output directory
STRUCTURE_FILE_NAME = "file_structure.json"
HASH_DB_FILE_NAME = "file_hashes.json"
DOT_REPLACEMENT = "__DOT__"

# Make sure this path is correct for your system
try:
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
except Exception:
    print("Warning: Tesseract path not found. OCR will fail if needed.")

def extract_text_from_pptx(pptx_stream):
    """
    Extracts text from a PowerPoint file using a MAXIMUM COMPLETENESS approach.
    1. Extracts native text from text boxes.
    2. Extracts images embedded in the slide, preprocesses them, and runs OCR.
    """
    print("  📽️ Starting PPTX text extraction (Native + OCR)...")
    all_slide_content = []
    
    try:
        # Presentation can take a file-like object (stream)
        prs = Presentation(pptx_stream)
        
        for i, slide in enumerate(prs.slides):
            print(f"  - Processing Slide {i+1}...")
            slide_text_parts = [f"--- Slide {i+1} ---"]
            
            # --- Part 1: Native Text Extraction ---
            native_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    native_text.append(shape.text)
            
            if native_text:
                joined_native = "\n".join(native_text)
                slide_text_parts.append(joined_native)
                print(f"    - Native text found: {len(joined_native)} chars.")

            # --- Part 2: OCR on Embedded Images ---
            ocr_text_list = []
            for shape in slide.shapes:
                # Check if the shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Extract image bytes
                        image_blob = shape.image.blob
                        img = Image.open(io.BytesIO(image_blob))
                        
                        # --- USE EXISTING PREPROCESSING ---
                        processed_img = preprocess_for_ocr(img)
                        
                        # Run Tesseract
                        text = pytesseract.image_to_string(processed_img)
                        
                        if text.strip():
                            ocr_text_list.append(text.strip())
                            
                    except Exception as ocr_err:
                        print(f"    - ⚠️ OCR warning on slide {i+1} image: {ocr_err}")

            if ocr_text_list:
                joined_ocr = "\n".join(ocr_text_list)
                print(f"    - OCR text found: {len(joined_ocr)} chars.")
                slide_text_parts.append("\n--- OCR Additions (Images) ---\n" + joined_ocr)

            # Combine all parts for this slide
            all_slide_content.append("\n".join(slide_text_parts))
        
        full_text = "\n\n".join(all_slide_content)
        print(f"  ✅ PPTX extraction complete. Total characters: {len(full_text)}")
        return full_text

    except Exception as e:
        print(f"  ❌ Error extracting PPTX: {e}")
        return ""

def preprocess_for_ocr(image: Image.Image) -> Image.Image:
    """
    Prepares an image for OCR by applying adaptive thresholding to preserve 
    colored text (like orange) that might be lost with global thresholding.
    """
    print("    - Pre-processing image for OCR (Adaptive Mode)...")
    # 1. Convert PIL Image to OpenCV format
    open_cv_image = np.array(image.convert('RGB'))
    open_cv_image = open_cv_image[:, :, ::-1].copy() # Convert RGB to BGR

    # 2. Grayscale
    gray = cv2.cvtColor(open_cv_image, cv2.COLOR_BGR2GRAY)

    # 3. Adaptive Thresholding
    # Instead of one global value, this calculates thresholds for small pixel areas.
    # This prevents light-colored text (orange) from being 'erased'.
    processed_image = cv2.adaptiveThreshold(
        gray, 
        255, 
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
        cv2.THRESH_BINARY, 
        11, # Block size (size of the local area)
        2   # Constant subtracted from the mean
    )

    # 4. Denoising (Optional)
    # Using a smaller kernel to avoid blurring thin text
    denoised = cv2.medianBlur(processed_image, 3)

    # 5. Convert back to PIL Image
    final_image = Image.fromarray(denoised)
    print("    - Pre-processing complete.")
    return final_image

def simplify_text(text):
    """A helper to normalize text for comparison by removing whitespace and making it lowercase."""
    return re.sub(r'\s+', '', text).lower()

def extract_text(pdf_stream):
    """
    Extracts text from a PDF stream using a MAXIMUM COMPLETENESS approach.
    For each page, it extracts native text and performs OCR, then intelligently
    merges the results to capture all possible content.

    Args:
        pdf_stream: A file-like object representing the PDF file.

    Returns:
        A string containing all extracted text from the PDF.
    """
    print("  🔎 Starting Maximum Completeness text extraction...")
    all_page_texts = []
    
    try:
        doc = fitz.open(stream=pdf_stream.read(), filetype="pdf")

        for page_num, page in enumerate(doc):
            print(f"  - Processing Page {page_num + 1}/{len(doc)}...")

            # --- Step 1: Get Native Text (The High-Quality Base) ---
            native_text = page.get_text("text")
            print(f"    - Native text found: {len(native_text)} chars.")

            # --- Step 2: Perform OCR (The Comprehensive Source) ---
            ocr_text = ""
            try:
                pix = page.get_pixmap(dpi=300)
                img_bytes = pix.tobytes("ppm")
                img = Image.open(io.BytesIO(img_bytes))

                # --- THIS IS THE NEW LINE ---
                processed_img = preprocess_for_ocr(img) 
                # ---------------------------

                ocr_text = pytesseract.image_to_string(processed_img) # Use the processed image
                print(f"    - OCR text found: {len(ocr_text)} chars.")
            except Exception as ocr_error:
                print(f"    - ❌ OCR failed for page {page_num + 1}: {ocr_error}")

            # --- Step 3: Intelligently Merge ---
            # If there's no native text, the page is purely an image. Use OCR text directly.
            if not native_text.strip():
                print("    - Verdict: Image-only page. Using OCR text.")
                all_page_texts.append(ocr_text)
                continue

            # If OCR text is negligible, the page is purely text. Use native text.
            if not ocr_text.strip():
                 print("    - Verdict: Text-only page. Using native text.")
                 all_page_texts.append(native_text)
                 continue

            # The complex case: Mixed content. Merge them.
            print("    - Verdict: Mixed content page. Merging results.")
            
            # Use the clean native text as our starting point.
            final_page_text = native_text
            
            # Create a simplified version of the native text for fast searching.
            simplified_native = simplify_text(native_text)

            # Find lines in OCR text that are NOT in the native text.
            unique_ocr_lines = []
            for line in ocr_text.splitlines():
                if line.strip() and simplify_text(line) not in simplified_native:
                    unique_ocr_lines.append(line)
            
            if unique_ocr_lines:
                print(f"    - Found {len(unique_ocr_lines)} unique lines from OCR. Appending them.")
                # Append the unique findings, separated clearly.
                unique_content = "\n".join(unique_ocr_lines)
                final_page_text += f"\n\n--- OCR Additions ---\n{unique_content}"

            all_page_texts.append(final_page_text)

    except Exception as e:
        print(f"  ❌ CRITICAL ERROR during PDF processing: {e}")
        return "\n\n".join(all_page_texts)
    
    finally:
        if 'doc' in locals() and doc:
            doc.close()
            
    full_text = "\n\n".join(all_page_texts)
    print(f"  ✅ Max-completeness extraction complete. Total characters: {len(full_text)}")
    return full_text
    

def extract_text_from_image(image_stream):
    """
    Extracts text from an image file stream using OCR.
    """
    print("  🖼️ Extracting text from image via OCR...")
    try:
        # 1. Open the image
        pil_image = Image.open(image_stream)
        
        # 2. Run our improved preprocessor
        processed_img = preprocess_for_ocr(pil_image)
        
        # 3. Perform OCR
        # --psm 3 tells Tesseract to look for blocks of text automatically
        custom_config = r'--oem 3 --psm 3'
        text = pytesseract.image_to_string(processed_img, config=custom_config)
        
        # 4. Fallback: If text is very short, try the raw grayscale version
        # (Sometimes colors are better read without thresholding)
        if len(text.strip()) < 10:
            print("    - Low text yield, trying grayscale fallback...")
            gray_img = pil_image.convert('L')
            text = pytesseract.image_to_string(gray_img, config=custom_config)

        print(f"  ✅ OCR extracted {len(text)} characters.")
        return text
        
    except Exception as e:
        print(f"  ❌ Image OCR failed: {e}")    
        return ""

def split_chunks(text):
    """Splits a large text into smaller chunks for easier processing."""
    print("  ✂️ Splitting text into chunks...")
    splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    chunks = splitter.split_text(text)
    print(f"  ✅ Created {len(chunks)} chunks.")
    return chunks

def delete_collection(coll_ref, batch_size):
    """Recursively deletes all documents and subcollections within a collection."""
    docs = coll_ref.limit(batch_size).stream()
    deleted = 0
    for doc in docs:
        for sub_coll_ref in doc.reference.collections():
            delete_collection(sub_coll_ref, batch_size)
        doc.reference.delete()
        deleted += 1
    if deleted >= batch_size:
        return delete_collection(coll_ref, batch_size)

def batch_save(collection, items, batch_size=100):
    batch = db.batch()
    for i, item in enumerate(items):
        if i % batch_size == 0 and i > 0:
            batch.commit()
            batch = db.batch()
        ref = collection.document()
        batch.set(ref, item)
    batch.commit()

def get_project_output_path(project_id: str) -> Path:
    """Returns the specific output path for a given project."""
    return TXT_OUTPUT_DIR / project_id

def init_txt_converter_for_project(project_id: str):
    project_output_path = get_project_output_path(project_id)
    project_output_path.mkdir(parents=True, exist_ok=True)
    # We only need to manage the structure file now
    structure_file = project_output_path / STRUCTURE_FILE_NAME
    if not structure_file.exists():
        structure_file.write_text("{}", encoding='utf-8')

def convert_to_txt(src_path: Path, txt_path: Path) -> bool:
    """Converts a source file to a .txt file with its content."""
    try:
        # For common text-based files, just copy the content directly.
        with open(src_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        txt_path.write_text(content, encoding='utf-8')
        return True
    except Exception as e:
        txt_path.write_text(f"[ERROR READING FILE: {src_path.name}]\n\n{e}\n\n[This might be a binary file or have an unsupported encoding.]", encoding='utf-8')
        return False

def build_file_tree(root_path: Path, allowed_extensions: list = None) -> dict:
    # ... (This function is correct, but ensure it creates paths to the FIRESTORE documents now)
    # Let's refine this function to be more generic.
    tree = {}
    if allowed_extensions:
        allowed_extensions = [f".{ext.lstrip('.')}" for ext in allowed_extensions]
    
    for item in sorted(root_path.rglob("*")):
        if item.is_file():
            if '.git' in item.parts: continue
            if allowed_extensions and item.suffix.lower() not in allowed_extensions: continue

            rel_path = item.relative_to(root_path)
            # We need the Firestore doc ID, which is a hash of the relative path.
            rel_path_str = str(rel_path).replace('\\', '/')
            doc_id = hashlib.sha1(rel_path_str.encode('utf-8')).hexdigest()

            parts = rel_path.parts
            d = tree
            for part in parts[:-1]:
                d = d.setdefault(part, {})
            # The value is now the Firestore document ID
            d[parts[-1]] = doc_id
    return tree

def load_hashes(project_id: str) -> dict:
    hash_db_file = get_project_output_path(project_id) / HASH_DB_FILE_NAME
    if hash_db_file.exists():
        return json.loads(hash_db_file.read_text(encoding='utf-8'))
    return {}

def save_hashes(project_id: str, hashes: dict):
    hash_db_file = get_project_output_path(project_id) / HASH_DB_FILE_NAME
    hash_db_file.write_text(json.dumps(hashes, indent=2), encoding='utf-8')

def is_git_repo(path: Path):
    try:
        return git.Repo(path, search_parent_directories=True)
    except (InvalidGitRepositoryError, NoSuchPathError):
        return None

def get_file_hash(filepath) -> str:
    """Computes SHA256 hash of a file's content."""
    hasher = hashlib.sha256()
    with open(filepath, 'rb') as f:
        while chunk := f.read(4096):
            hasher.update(chunk)
    return hasher.hexdigest()

def get_converted_file_ref(db, project_id, original_path_str: str, sub_collection: str, top_level_collection: str = "projects"):
    """Creates a consistent and safe document ID from the original file path."""
    path_hash = hashlib.sha1(original_path_str.encode('utf-8')).hexdigest()
    return db.collection(top_level_collection).document(project_id).collection(sub_collection).document(path_hash)

def convert_and_upload_to_firestore(db, project_id, file_path, source_root, sub_collection: str, top_level_collection: str):
    """
    Reads a local file → text → uploads to Firestore.
    Accepts top-level and sub-collection names for flexibility.
    Returns (hash, doc_id) on success, otherwise None.
    """
    rel_path_str = str(file_path.relative_to(source_root)).replace('\\', '/')
    print(f"  Processing: {rel_path_str}")

    try:
        content = file_path.read_text(encoding='utf-8', errors='ignore')
        current_hash = hashlib.sha256(content.encode('utf-8')).hexdigest()

        # --- MODIFIED: Use the provided top_level_collection ---
        doc_ref = db.collection(top_level_collection) \
                    .document(project_id) \
                    .collection(sub_collection) \
                    .document()

        doc_ref.set({
            'original_path': rel_path_str,
            'content': content,
            'hash': current_hash,
            'timestamp': firestore.SERVER_TIMESTAMP,
        })

        print(f"    -> Uploaded to '{top_level_collection}/{project_id}/{sub_collection}' (doc_id={doc_ref.id})")
        return current_hash, doc_ref.id

    except Exception as e:
        print(f"    -> FAILED {rel_path_str}: {e}")
        return None

# cache
class SimpleL1Cache:
    def __init__(self, max_size=256, ttl=10):
        """
        A simple in-memory LRU (Least Recently Used) cache.
        :param max_size: The maximum number of items to store.
        :param ttl: Time-to-live in seconds for each item.
        """
        self.cache = OrderedDict()
        self.max_size = max_size
        self.ttl = ttl

    def get(self, key):
        if key not in self.cache:
            return None
        
        value, expiry = self.cache[key]
        
        # Check if the item has expired
        if time.time() > expiry:
            del self.cache[key]
            return None
            
        # Move the item to the end to mark it as recently used
        self.cache.move_to_end(key)
        return value

    def set(self, key, value):
        # Check if we need to make space
        if len(self.cache) >= self.max_size:
            # Remove the oldest item
            self.cache.popitem(last=False)
            
        expiry = time.time() + self.ttl
        self.cache[key] = (value, expiry)

L1_CACHE = SimpleL1Cache(max_size=512, ttl=20)

# Add these functions to utils.py

def generate_tree_text_from_paths(root_name: str, file_paths: list) -> str:
    """
    Industrial-grade Tree Generator using Trie traversal.
    Ensures perfect vertical bar continuity and trailing slash folder indicators.
    """
    # 1. Build a Trie (Prefix Tree) from the flat list of paths
    trie = {}
    for path_str in file_paths:
        # Use Path.parts to handle both \ and / correctly
        parts = Path(path_str).parts
        current_level = trie
        for part in parts:
            if part not in current_level:
                current_level[part] = {}
            current_level = current_level[part]

    # 2. Recursive walk with "Last-Sibling" detection
    def walk_trie(node, prefix=""):
        lines = []
        # Sort keys to maintain alphabetical order (standard tree behavior)
        items = sorted(node.keys())
        
        for i, name in enumerate(items):
            is_last = (i == len(items) - 1)
            
            # Use standard box-drawing characters
            connector = "└── " if is_last else "├── "
            
            # Folders are nodes that have children in the Trie
            is_dir = len(node[name]) > 0
            display_name = f"{name}/" if is_dir else name
            
            lines.append(f"{prefix}{connector}{display_name}")
            
            # If it's a directory, recurse into children
            if is_dir:
                # If this is the last sibling, the children's prefix is empty space.
                # If there are more siblings below, the children's prefix needs a vertical bar.
                extension = "    " if is_last else "│   "
                lines.extend(walk_trie(node[name], prefix + extension))
        
        return lines

    # 3. Generate final string
    tree_lines = walk_trie(trie)
    return f"{root_name}/\n" + "\n".join(tree_lines)

def generate_tree_with_stats(root_name: str, file_paths: list, files_metadata: dict = None) -> str:
    """
    Enhanced tree generation that includes file statistics.
    
    Args:
        root_name: The name of the root directory
        file_paths: List of relative file paths
        files_metadata: Optional dict mapping paths to metadata (size, lines, etc.)
    
    Returns:
        A formatted tree structure with statistics
    """
    tree = {}
    file_count = 0
    dir_count = 0
    total_size = 0
    
    for path_str in file_paths:
        parts = Path(path_str).parts
        current_dict = tree
        
        for part in parts[:-1]:
            if part not in current_dict:
                dir_count += 1
            current_dict = current_dict.setdefault(part, {})
        
        # Store file with metadata
        file_info = {'_is_file': True}
        if files_metadata and path_str in files_metadata:
            file_info.update(files_metadata[path_str])
            total_size += files_metadata[path_str].get('size', 0)
        
        current_dict[parts[-1]] = file_info
        file_count += 1

    def build_tree_lines(sub_tree, prefix=""):
        lines = []
        items = sorted(sub_tree.keys())
        
        for i, key in enumerate(items):
            is_last = (i == len(items) - 1)
            value = sub_tree[key]
            connector = "└── " if is_last else "├── "
            
            if isinstance(value, dict) and value.get('_is_file'):
                # It's a file with metadata
                size_info = ""
                if 'size' in value:
                    size_kb = value['size'] / 1024
                    size_info = f" ({size_kb:.1f} KB)"
                lines.append(f"{prefix}{connector}{key}{size_info}")
            elif isinstance(value, dict):
                # It's a directory
                lines.append(f"{prefix}{connector}{key}/")
                extension = "    " if is_last else "│   "
                new_prefix = prefix + extension
                lines.extend(build_tree_lines(value, new_prefix))
            else:
                # Simple file
                lines.append(f"{prefix}{connector}{key}")
                
        return lines

    tree_lines = build_tree_lines(tree)
    
    # Create summary header
    summary = f"""
    {root_name}/
    {'='*60}
    Directories: {dir_count}
    Files: {file_count}
    Total Size: {total_size / 1024 / 1024:.2f} MB
    {'='*60}
    """
    
    full_tree = summary + "\n".join(tree_lines)
    return full_tree


def format_bytes(size_bytes: int) -> str:
    """
    Converts bytes to human-readable format.
    
    Args:
        size_bytes: Size in bytes
    
    Returns:
        Formatted string (e.g., "1.5 MB", "348 KB")
    """
    for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.1f} PB"


def validate_tree_structure(tree: dict) -> bool:
    """
    Validates that a tree structure is properly formed.
    
    Args:
        tree: Dictionary representing the file tree
    
    Returns:
        True if valid, False otherwise
    """
    def check_node(node):
        if not isinstance(node, dict):
            return False
        
        for key, value in node.items():
            if not isinstance(key, str):
                return False
            
            # Files are marked with None or have _is_file flag
            if value is None or (isinstance(value, dict) and value.get('_is_file')):
                continue
            
            # Directories must be dicts
            if not isinstance(value, dict):
                return False
            
            # Recursively check subdirectories
            if not check_node(value):
                return False
        
        return True
    
    return check_node(tree)